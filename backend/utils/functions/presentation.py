from pptx import Presentation
import os
from pptx.util import Inches, Pt

def generate_presentation(presentation_json, file_url):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    
    for json_slide in presentation_json:
        slide = prs.slides.add_slide(layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)

        title_shape.text = json_slide["TITLE"]


        if json_slide["TEXT"]["TYPE"] == "BULLETS":
            tf = body_shape.text_frame

            bullets = json_slide["TEXT"]["VALUE"]
            for bullet in bullets:
                tf.text = bullet
        else:
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = json_slide["TEXT"]["VALUE"]
            p.font.size = Pt(40)

            
    prs.save(file_url)